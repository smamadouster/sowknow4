"""
Text extraction service for various document formats
"""

import logging
import subprocess
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)


class TextExtractor:
    """Service for extracting text from various document formats"""

    def __init__(self) -> None:
        self.supported_formats = {
            ".pdf": self._extract_from_pdf,
            ".docx": self._extract_from_docx,
            ".doc": self._extract_from_doc,
            ".pptx": self._extract_from_pptx,
            ".ppsx": self._extract_from_pptx,
            ".ppt": self._extract_from_ppt,
            ".xlsx": self._extract_from_xlsx,
            ".xls": self._extract_from_xls,
            ".xlt": self._extract_from_xls,
            ".xltx": self._extract_from_xlsx,
            ".txt": self._extract_from_txt,
            ".md": self._extract_from_txt,
            ".json": self._extract_from_json,
            ".csv": self._extract_from_csv,
            ".xml": self._extract_from_xml,
            ".html": self._extract_from_html,
            ".htm": self._extract_from_html,
            ".epub": self._extract_from_epub,
            ".rtf": self._extract_from_rtf,
            ".zip": self._extract_from_zip,
            ".xmind": self._extract_from_zip,
            ".msg": self._extract_from_msg,
            ".oft": self._extract_from_msg,
        }

    def get_file_extension(self, filename: str) -> str:
        """Get file extension from filename"""
        return Path(filename).suffix.lower()

    async def extract_text(self, file_path: str, filename: str) -> dict[str, Any]:
        """
        Extract text from a document file

        Args:
            file_path: Path to the file
            filename: Original filename

        Returns:
            dict with extracted text, metadata, and page count
        """
        file_extension = self.get_file_extension(filename)

        if file_extension not in self.supported_formats:
            return {
                "text": "",
                "error": f"Unsupported file format: {file_extension}",
                "pages": 0,
            }

        try:
            extractor = self.supported_formats[file_extension]
            result = await extractor(file_path)

            # Add metadata
            result["format"] = file_extension
            result["success"] = "error" not in result

            return result

        except Exception as e:
            logger.error(f"Error extracting text from {filename}: {str(e)}")
            return {
                "text": "",
                "error": str(e),
                "format": file_extension,
                "pages": 0,
                "success": False,
            }

    async def _extract_from_pdf(self, file_path: str) -> dict[str, Any]:
        """Extract text from PDF file"""
        try:
            import PyPDF2

            text_parts = []
            with open(file_path, "rb") as file:
                pdf_reader = PyPDF2.PdfReader(file)
                page_count = len(pdf_reader.pages)

                for page_num, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text()
                    if page_text.strip():
                        text_parts.append(f"[Page {page_num + 1}]\n{page_text}")

            return {
                "text": "\n\n".join(text_parts),
                "pages": page_count,
                "source": "pypdf2",
            }

        except Exception as e:
            logger.error(f"Error extracting from PDF: {str(e)}")
            return {"text": "", "error": str(e), "pages": 0}

    async def _extract_from_docx(self, file_path: str) -> dict[str, Any]:
        """Extract text from DOCX file"""
        try:
            from docx import Document

            doc = Document(file_path)
            text_parts = []

            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)

            # Extract from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip() for cell in row.cells])
                    if row_text.strip():
                        text_parts.append(row_text)

            return {"text": "\n".join(text_parts), "pages": 0, "source": "python-docx"}

        except Exception as e:
            logger.error(f"Error extracting from DOCX: {str(e)}")
            return {"text": "", "error": str(e), "pages": 0}

    async def _extract_from_doc(self, file_path: str) -> dict[str, Any]:
        """Extract text from legacy DOC file using antiword, with python-docx fallback"""
        # Try antiword first (handles binary .doc format)
        try:
            result = subprocess.run(  # noqa: S603
                ["antiword", "-w", "0", file_path],
                capture_output=True,
                text=True,
                timeout=60,
            )
            if result.returncode == 0 and result.stdout.strip():
                return {"text": result.stdout.strip(), "pages": 0, "source": "antiword"}
            if result.stderr:
                logger.warning(f"antiword stderr for {file_path}: {result.stderr.strip()}")
        except FileNotFoundError:
            logger.warning("antiword not installed, trying python-docx fallback")
        except subprocess.TimeoutExpired:
            logger.warning(f"antiword timed out for {file_path}")
        except Exception as e:
            logger.warning(f"antiword failed for {file_path}: {e}")

        # Fallback: some .doc files are actually OOXML and python-docx can read them
        try:
            from docx import Document
            doc = Document(file_path)
            text_parts = [p.text for p in doc.paragraphs if p.text.strip()]
            for table in doc.tables:
                for row in table.rows:
                    row_text = "\t".join(cell.text for cell in row.cells)
                    if row_text.strip():
                        text_parts.append(row_text)
            if text_parts:
                return {"text": "\n".join(text_parts), "pages": 0, "source": "python-docx-fallback"}
        except Exception as e:
            logger.warning(f"python-docx fallback failed for {file_path}: {e}")

        return {
            "text": "",
            "error": "Could not extract text from DOC file. Ensure antiword is installed.",
            "pages": 0,
            "source": "error",
        }

    async def _extract_from_pptx(self, file_path: str) -> dict[str, Any]:
        """Extract text from PPTX file"""
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            text_parts = []
            slide_count = len(prs.slides)

            for slide_num, slide in enumerate(prs.slides):
                slide_text_parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text_parts.append(shape.text)

                if slide_text_parts:
                    text_parts.append(f"[Slide {slide_num + 1}]\n" + "\n".join(slide_text_parts))

            return {
                "text": "\n\n".join(text_parts),
                "pages": slide_count,
                "source": "python-pptx",
            }

        except Exception as e:
            logger.error(f"Error extracting from PPTX: {str(e)}")
            return {"text": "", "error": str(e), "pages": 0}

    async def _extract_from_ppt(self, file_path: str) -> dict[str, Any]:
        """Extract text from legacy PPT file using catppt."""
        try:
            import subprocess

            result = subprocess.run(
                ["catppt", file_path],
                capture_output=True,
                text=True,
                timeout=60,
            )
            if result.returncode != 0:
                return {
                    "text": "",
                    "error": f"catppt failed: {result.stderr}",
                    "pages": 0,
                    "source": "catppt-error",
                }

            text = result.stdout.strip()
            # catppt emits page breaks as "^L" (form feed); count them as pages
            pages = text.count("\f") + 1 if text else 0
            return {
                "text": text,
                "pages": pages,
                "source": "catppt",
            }
        except FileNotFoundError:
            return {
                "text": "",
                "error": "Legacy PPT files require catppt (install catdoc package)",
                "pages": 0,
                "source": "error",
            }
        except Exception as e:
            return {
                "text": "",
                "error": f"PPT extraction failed: {e}",
                "pages": 0,
                "source": "error",
            }

    async def _extract_from_xlsx(self, file_path: str) -> dict[str, Any]:
        """Extract text from XLSX file.

        Rows are batched into blocks of ~50 rows so the chunking service can
        create proper 512-token chunks instead of one chunk per row (which
        caused 40k+ chunks from a single spreadsheet).

        Safety bounds:
        - Max 20 visible sheets (skip hidden/very-hidden)
        - Max 5,000 rows per sheet
        - Max 500,000 characters total text (truncate with notice)
        - Skip rows where all cells are empty/whitespace
        """
        try:
            from openpyxl import load_workbook

            _MAX_SHEETS = 20
            _MAX_ROWS_PER_SHEET = 5000
            _MAX_TOTAL_CHARS = 500_000
            _ROW_BATCH_SIZE = 50

            wb = load_workbook(file_path, read_only=True)
            text_parts = []
            sheets_processed = 0
            total_chars = 0
            truncated = False

            for sheet_name in wb.sheetnames:
                if sheets_processed >= _MAX_SHEETS:
                    break

                sheet = wb[sheet_name]
                # Skip hidden/very-hidden sheets
                if sheet.sheet_state != "visible":
                    continue

                sheets_processed += 1
                sheet_text_parts = [f"[Sheet: {sheet_name}]"]
                batch_rows: list[str] = []
                row_count = 0

                for row in sheet.iter_rows(values_only=True):
                    row_count += 1
                    if row_count > _MAX_ROWS_PER_SHEET:
                        sheet_text_parts.append("[... rows truncated at 5000 ...]")
                        break

                    # Skip rows where all cells are None or whitespace-only
                    if not any(cell is not None and str(cell).strip() for cell in row):
                        continue

                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    batch_rows.append(row_text)

                    if len(batch_rows) >= _ROW_BATCH_SIZE:
                        batch_text = "\n".join(batch_rows)
                        if total_chars + len(batch_text) > _MAX_TOTAL_CHARS:
                            truncated = True
                            break
                        sheet_text_parts.append(batch_text)
                        total_chars += len(batch_text) + 2  # +2 for \n\n
                        batch_rows = []

                if not truncated and batch_rows:
                    batch_text = "\n".join(batch_rows)
                    if total_chars + len(batch_text) <= _MAX_TOTAL_CHARS:
                        sheet_text_parts.append(batch_text)
                        total_chars += len(batch_text) + 2
                    else:
                        truncated = True

                if len(sheet_text_parts) > 1:
                    text_parts.append("\n\n".join(sheet_text_parts))

                if truncated:
                    text_parts.append("[... document truncated at 500,000 characters ...]")
                    break

            return {
                "text": "\n\n".join(text_parts),
                "pages": sheets_processed,
                "source": "openpyxl",
            }

        except Exception as e:
            logger.error(f"Error extracting from XLSX: {str(e)}")
            return {"text": "", "error": str(e), "pages": 0}

    async def _extract_from_xls(self, file_path: str) -> dict[str, Any]:
        """Extract text from legacy XLS file using xlrd.

        Uses the same safety bounds as XLSX to prevent monster documents.
        If xlrd fails (e.g. missing dependency or newer-format file masquerading
        as .xlt), falls back to openpyxl.
        """
        try:
            import xlrd

            _MAX_SHEETS = 20
            _MAX_ROWS_PER_SHEET = 5000
            _MAX_TOTAL_CHARS = 500_000
            _ROW_BATCH_SIZE = 50

            wb = xlrd.open_workbook(file_path, on_demand=True)
            text_parts = []
            sheets_processed = 0
            total_chars = 0
            truncated = False

            for sheet_idx in range(min(wb.nsheets, _MAX_SHEETS)):
                sheet = wb.sheet_by_index(sheet_idx)
                sheets_processed += 1
                sheet_text_parts = [f"[Sheet: {sheet.name}]"]
                batch_rows: list[str] = []
                row_count = min(sheet.nrows, _MAX_ROWS_PER_SHEET)

                if sheet.nrows > _MAX_ROWS_PER_SHEET:
                    row_count = _MAX_ROWS_PER_SHEET

                for row_idx in range(row_count):
                    if row_idx >= _MAX_ROWS_PER_SHEET:
                        sheet_text_parts.append("[... rows truncated at 5000 ...]")
                        break

                    row_values = sheet.row_values(row_idx)
                    # Skip rows where all cells are empty/whitespace
                    if not any(cell is not None and str(cell).strip() for cell in row_values):
                        continue

                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row_values])
                    batch_rows.append(row_text)

                    if len(batch_rows) >= _ROW_BATCH_SIZE:
                        batch_text = "\n".join(batch_rows)
                        if total_chars + len(batch_text) > _MAX_TOTAL_CHARS:
                            truncated = True
                            break
                        sheet_text_parts.append(batch_text)
                        total_chars += len(batch_text) + 2
                        batch_rows = []

                if not truncated and batch_rows:
                    batch_text = "\n".join(batch_rows)
                    if total_chars + len(batch_text) <= _MAX_TOTAL_CHARS:
                        sheet_text_parts.append(batch_text)
                        total_chars += len(batch_text) + 2
                    else:
                        truncated = True

                if len(sheet_text_parts) > 1:
                    text_parts.append("\n\n".join(sheet_text_parts))

                if truncated:
                    text_parts.append("[... document truncated at 500,000 characters ...]")
                    break

            wb.release_resources()

            return {
                "text": "\n\n".join(text_parts),
                "pages": sheets_processed,
                "source": "xlrd",
            }

        except Exception as e:
            logger.warning(f"xlrd failed for {file_path}: {e} — trying openpyxl fallback")
            # Some .xlt files are actually newer-format templates; try openpyxl
            try:
                fallback = await self._extract_from_xlsx(file_path)
                if fallback.get("text") or not fallback.get("error"):
                    fallback["source"] = "openpyxl-fallback"
                    return fallback
            except Exception as fallback_err:
                logger.warning(f"openpyxl fallback also failed for {file_path}: {fallback_err}")

            logger.error(f"Error extracting from XLS: {str(e)}")
            return {"text": "", "error": str(e), "pages": 0}

    async def _extract_from_txt(self, file_path: str) -> dict[str, Any]:
        """Extract text from TXT or MD file"""
        try:
            # Try UTF-8 first, fallback to latin-1
            try:
                with open(file_path, encoding="utf-8") as f:
                    text = f.read()
            except UnicodeDecodeError:
                with open(file_path, encoding="latin-1") as f:
                    text = f.read()

            return {"text": text, "pages": 0, "source": "file"}

        except Exception as e:
            logger.error(f"Error extracting from TXT: {str(e)}")
            return {"text": "", "error": str(e), "pages": 0}

    async def _extract_from_json(self, file_path: str) -> dict[str, Any]:
        """Extract text from JSON file"""
        try:
            import json

            with open(file_path, encoding="utf-8") as f:
                data = json.load(f)

            # Convert JSON to readable text
            text = json.dumps(data, indent=2, ensure_ascii=False)

            return {"text": text, "pages": 0, "source": "json"}

        except Exception as e:
            logger.error(f"Error extracting from JSON: {str(e)}")
            return {"text": "", "error": str(e), "pages": 0}

    async def _extract_from_csv(self, file_path: str) -> dict[str, Any]:
        """Extract text from CSV file"""
        try:
            import csv

            try:
                with open(file_path, encoding="utf-8", newline="") as f:
                    reader = csv.reader(f)
                    rows = list(reader)
            except UnicodeDecodeError:
                with open(file_path, encoding="latin-1", newline="") as f:
                    reader = csv.reader(f)
                    rows = list(reader)

            text_lines = []
            for row in rows:
                line = " | ".join(str(cell) for cell in row)
                if line.strip():
                    text_lines.append(line)

            return {
                "text": "\n".join(text_lines),
                "pages": 0,
                "source": "csv",
            }

        except Exception as e:
            logger.error(f"Error extracting from CSV: {str(e)}")
            return {"text": "", "error": str(e), "pages": 0}

    async def _extract_from_xml(self, file_path: str) -> dict[str, Any]:
        """Extract text from XML file"""
        try:
            import xml.etree.ElementTree as ET

            tree = ET.parse(file_path)  # noqa: S314
            root = tree.getroot()

            def extract_text_recursive(element: ET.Element) -> list[str]:
                parts = []
                if element.text and element.text.strip():
                    parts.append(element.text.strip())
                for child in element:
                    parts.extend(extract_text_recursive(child))
                    if child.tail and child.tail.strip():
                        parts.append(child.tail.strip())
                return parts

            text_parts = extract_text_recursive(root)

            return {
                "text": "\n".join(text_parts),
                "pages": 0,
                "source": "xml",
            }

        except Exception as e:
            logger.error(f"Error extracting from XML: {str(e)}")
            return {"text": "", "error": str(e), "pages": 0}

    async def _extract_from_epub(self, file_path: str) -> dict[str, Any]:
        """Extract text from EPUB file"""
        try:
            from ebooklib import epub

            book = epub.read_epub(file_path)
            text_parts = []

            for item in book.get_items():
                if item.get_type() == epub.EpubHtml:
                    # Remove HTML tags
                    import re

                    content = item.get_content().decode("utf-8")
                    text = re.sub(r"<[^>]+>", "", content)
                    text = " ".join(text.split())
                    if text.strip():
                        text_parts.append(text)

            return {"text": "\n\n".join(text_parts), "pages": 0, "source": "ebooklib"}

        except Exception as e:
            logger.error(f"Error extracting from EPUB: {str(e)}")
            return {"text": "", "error": str(e), "pages": 0}

    async def _extract_from_html(self, file_path: str) -> dict[str, Any]:
        """Extract readable text from HTML/HTM files."""
        try:
            import html
            import re

            try:
                with open(file_path, encoding="utf-8") as f:
                    data = f.read()
            except UnicodeDecodeError:
                with open(file_path, encoding="latin-1") as f:
                    data = f.read()

            data = re.sub(r"(?is)<(script|style).*?</\1>", " ", data)
            text = re.sub(r"(?s)<[^>]+>", " ", data)
            text = html.unescape(text)
            text = re.sub(r"[ \t\r\f\v]+", " ", text)
            text = re.sub(r"\n\s+", "\n", text)
            return {"text": text.strip(), "pages": 0, "source": "html-parser"}

        except Exception as e:
            logger.error(f"Error extracting from HTML: {str(e)}")
            return {"text": "", "error": str(e), "pages": 0}

    async def _extract_from_rtf(self, file_path: str) -> dict[str, Any]:
        """Extract readable text from RTF using a lightweight parser."""
        try:
            import re

            try:
                with open(file_path, encoding="utf-8") as f:
                    data = f.read()
            except UnicodeDecodeError:
                with open(file_path, encoding="latin-1") as f:
                    data = f.read()

            def replace_hex(match: re.Match[str]) -> str:
                try:
                    return bytes.fromhex(match.group(1)).decode("latin-1")
                except Exception:
                    return " "

            text = re.sub(r"\\'([0-9a-fA-F]{2})", replace_hex, data)
            text = re.sub(r"\\par[d]?", "\n", text)
            text = re.sub(r"\\tab", "\t", text)
            text = re.sub(r"\\[a-zA-Z]+-?\d* ?", " ", text)
            text = re.sub(r"[{}]", " ", text)
            text = re.sub(r"[ \t]+", " ", text)
            text = re.sub(r"\n\s+", "\n", text)
            return {"text": text.strip(), "pages": 0, "source": "rtf-parser"}

        except Exception as e:
            logger.error(f"Error extracting from RTF: {str(e)}")
            return {"text": "", "error": str(e), "pages": 0}

    async def _extract_from_zip(self, file_path: str) -> dict[str, Any]:
        """Index ZIP archive contents as a searchable manifest."""
        try:
            import zipfile

            lines = ["ZIP archive contents:"]
            with zipfile.ZipFile(file_path) as archive:
                infos = archive.infolist()
                for info in infos:
                    if info.is_dir():
                        continue
                    lines.append(f"{info.filename} ({info.file_size} bytes)")

            return {
                "text": "\n".join(lines),
                "pages": 0,
                "source": "zip-manifest",
                "file_count": max(len(lines) - 1, 0),
            }

        except Exception as e:
            logger.error(f"Error extracting from ZIP: {str(e)}")
            return {"text": "", "error": str(e), "pages": 0}

    async def _extract_from_msg(self, file_path: str) -> dict[str, Any]:
        """Extract searchable strings from Outlook MSG files without extra dependencies."""
        try:
            import re

            with open(file_path, "rb") as f:
                data = f.read()

            parts: list[str] = []
            for match in re.finditer(rb"(?:[\x20-\x7e]\x00){6,}", data):
                try:
                    text = match.group(0).decode("utf-16le", errors="ignore").strip()
                    if text:
                        parts.append(text)
                except Exception:
                    pass

            for match in re.finditer(rb"[\x20-\x7e]{8,}", data):
                try:
                    text = match.group(0).decode("latin-1", errors="ignore").strip()
                    if text:
                        parts.append(text)
                except Exception:
                    pass

            seen: set[str] = set()
            unique_parts: list[str] = []
            for part in parts:
                cleaned = " ".join(part.split())
                if cleaned and cleaned not in seen:
                    seen.add(cleaned)
                    unique_parts.append(cleaned)

            return {
                "text": "\n".join(unique_parts),
                "pages": 0,
                "source": "msg-strings",
            }

        except Exception as e:
            logger.error(f"Error extracting from MSG: {str(e)}")
            return {"text": "", "error": str(e), "pages": 0}

    async def extract_images_from_pdf(self, file_path: str) -> list[bytes]:
        """
        Extract images from PDF for OCR processing

        Args:
            file_path: Path to PDF file

        Returns:
            List of image bytes
        """
        try:
            import PyPDF2

            images = []
            with open(file_path, "rb") as file:
                pdf_reader = PyPDF2.PdfReader(file)

                for _page_num, page in enumerate(pdf_reader.pages):
                    if "/XObject" in page["/Resources"]:
                        xObject = page["/Resources"]["/XObject"].get_object()

                        for obj_name in xObject:
                            if xObject[obj_name]["/Subtype"] == "/Image":
                                try:
                                    image_data = xObject[obj_name]._data
                                    images.append(image_data)
                                except (KeyError, AttributeError, IndexError):
                                    continue

            return images

        except Exception as e:
            logger.error(f"Error extracting images from PDF: {str(e)}")
            return []


# Global text extractor instance
text_extractor = TextExtractor()
